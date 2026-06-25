"""
slide_tools.py - Slide Agent용 PowerPoint 생성 도구
"""
import json, os
from pathlib import Path
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

OUTPUT_DIR = Path(os.getenv("SLIDE_OUTPUT_DIR", "data/slides"))


def _check():
    if not PPTX_AVAILABLE:
        return {"success": False, "error": "python-pptx 미설치. pip install python-pptx 실행 필요"}
    return None


def create_presentation(filename: str, title: str, subtitle: str = "") -> dict:
    """새 프레젠테이션을 생성하고 타이틀 슬라이드를 추가한다."""
    err = _check()
    if err:
        return err
    try:
        OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # 타이틀 슬라이드
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        if subtitle and slide.placeholders[1]:
            slide.placeholders[1].text = subtitle

        if not filename.endswith(".pptx"):
            filename += ".pptx"
        path = OUTPUT_DIR / filename
        prs.save(str(path))
        return {"success": True, "path": str(path), "message": f"프레젠테이션 생성: {filename}"}
    except Exception as e:
        return {"success": False, "error": str(e)}


def add_content_slide(pptx_path: str, title: str, bullet_points: list, notes: str = "") -> dict:
    """불릿 포인트가 있는 콘텐츠 슬라이드를 추가한다."""
    err = _check()
    if err:
        return err
    try:
        prs = Presentation(pptx_path)
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, point in enumerate(bullet_points):
            if i == 0:
                tf.paragraphs[0].text = point
            else:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        prs.save(pptx_path)
        slide_num = len(prs.slides)
        return {"success": True, "message": f"슬라이드 {slide_num} 추가: {title}"}
    except Exception as e:
        return {"success": False, "error": str(e)}


def add_section_slide(pptx_path: str, section_title: str) -> dict:
    """섹션 구분 슬라이드를 추가한다."""
    err = _check()
    if err:
        return err
    try:
        prs = Presentation(pptx_path)
        slide = prs.slides.add_slide(prs.slide_layouts[2])
        slide.shapes.title.text = section_title
        prs.save(pptx_path)
        return {"success": True, "message": f"섹션 슬라이드 추가: {section_title}"}
    except Exception as e:
        return {"success": False, "error": str(e)}


def list_slides(pptx_path: str) -> dict:
    """프레젠테이션의 슬라이드 목록을 반환한다."""
    err = _check()
    if err:
        return err
    try:
        prs = Presentation(pptx_path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            title = ""
            if slide.shapes.title:
                title = slide.shapes.title.text
            slides.append({"index": i, "title": title})
        return {"success": True, "total": len(slides), "slides": slides, "path": pptx_path}
    except Exception as e:
        return {"success": False, "error": str(e)}


def list_presentations() -> dict:
    """저장된 프레젠테이션 목록을 반환한다."""
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    files = [{"name": f.name, "path": str(f), "size_kb": round(f.stat().st_size / 1024, 1)}
             for f in OUTPUT_DIR.glob("*.pptx")]
    return {"success": True, "count": len(files), "files": files}


SLIDE_TOOLS = [
    {
        "name": "create_presentation",
        "description": "새 PowerPoint 프레젠테이션을 생성하고 타이틀 슬라이드를 추가한다.",
        "input_schema": {
            "type": "object",
            "properties": {
                "filename": {"type": "string", "description": "저장할 파일명 (예: my_presentation)"},
                "title":    {"type": "string", "description": "프레젠테이션 제목"},
                "subtitle": {"type": "string", "description": "부제목 (선택)"}
            },
            "required": ["filename", "title"]
        }
    },
    {
        "name": "add_content_slide",
        "description": "제목과 불릿 포인트가 있는 콘텐츠 슬라이드를 추가한다.",
        "input_schema": {
            "type": "object",
            "properties": {
                "pptx_path":    {"type": "string", "description": "pptx 파일 경로"},
                "title":        {"type": "string", "description": "슬라이드 제목"},
                "bullet_points":{"type": "array", "items": {"type": "string"}, "description": "불릿 포인트 목록"},
                "notes":        {"type": "string", "description": "발표자 노트 (선택)"}
            },
            "required": ["pptx_path", "title", "bullet_points"]
        }
    },
    {
        "name": "add_section_slide",
        "description": "섹션 구분 슬라이드를 추가한다.",
        "input_schema": {
            "type": "object",
            "properties": {
                "pptx_path":     {"type": "string", "description": "pptx 파일 경로"},
                "section_title": {"type": "string", "description": "섹션 제목"}
            },
            "required": ["pptx_path", "section_title"]
        }
    },
    {
        "name": "list_slides",
        "description": "프레젠테이션의 슬라이드 목록을 조회한다.",
        "input_schema": {
            "type": "object",
            "properties": {
                "pptx_path": {"type": "string", "description": "pptx 파일 경로"}
            },
            "required": ["pptx_path"]
        }
    },
    {
        "name": "list_presentations",
        "description": "저장된 프레젠테이션 파일 목록을 조회한다.",
        "input_schema": {"type": "object", "properties": {}, "required": []}
    }
]


def execute_tool(tool_name: str, tool_input: dict) -> str:
    tool_map = {
        "create_presentation": create_presentation,
        "add_content_slide":   add_content_slide,
        "add_section_slide":   add_section_slide,
        "list_slides":         list_slides,
        "list_presentations":  list_presentations,
    }
    if tool_name not in tool_map:
        return json.dumps({"success": False, "error": f"Unknown tool: {tool_name}"})
    return json.dumps(tool_map[tool_name](**tool_input), ensure_ascii=False)
